#!/usr/bin/env python3
"""Generate EUKOV Management Portal presentation (.pptx) — aligned with project PDF + codebase."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "EUKOV_Management_Portal.pptx"

BG = RGBColor(0xFA, 0xF7, 0xF2)
ACCENT = RGBColor(0xC4, 0x6D, 0x3A)
DARK = RGBColor(0x1E, 0x2A, 0x3A)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BOX_FILL = RGBColor(0xFF, 0xFD, 0xF8)
BOX_BORDER = RGBColor(0xE8, 0xD8, 0xC4)
AI_FILL = RGBColor(0xF3, 0xE6, 0xD4)
TABLE_HEADER = RGBColor(0xF5, 0xED, 0xE3)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=0.42, height=0.55) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(top), Inches(0.08), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_slide_title(slide, title: str, top=0.38) -> None:
    add_accent_bar(slide, top=top)
    box = slide.shapes.add_textbox(Inches(0.78), Inches(top - 0.05), Inches(11.5), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Georgia"


def add_center_title_slide(slide, title: str, subtitle: str, footer: str | None = None) -> None:
    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.5), Inches(-0.5), Inches(4), Inches(4)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xC4)
    glow.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(11.7), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run1 = p.add_run()
    run1.text = "EUKOV "
    run1.font.size = Pt(44)
    run1.font.bold = True
    run1.font.color.rgb = DARK
    run1.font.name = "Georgia"
    run2 = p.add_run()
    run2.text = "Management Portal"
    run2.font.size = Pt(44)
    run2.font.bold = True
    run2.font.color.rgb = ACCENT
    run2.font.name = "Georgia"

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(17)
    p2.font.color.rgb = MUTED
    p2.font.name = "Calibri"
    p2.space_before = Pt(16)

    if footer:
        p3 = tf.add_paragraph()
        p3.text = footer
        p3.alignment = PP_ALIGN.CENTER
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(0x9A, 0x9A, 0x9A)
        p3.font.name = "Calibri"
        p3.space_before = Pt(28)


def add_two_column_boxes(slide, left_title, left_body, right_title, right_body, top=1.45) -> None:
    for col, (title, body) in enumerate([(left_title, left_body), (right_title, right_body)]):
        left = Inches(0.75 + col * 6.15)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(top), Inches(5.85), Inches(4.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = BOX_BORDER
        shape.line.width = Pt(1)
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, Inches(top), Inches(5.85), Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()

        tf = shape.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.35)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.font.name = "Georgia"
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(15)
        p2.font.color.rgb = MUTED
        p2.font.name = "Calibri"
        p2.space_before = Pt(12)


def add_bullet_list(slide, items: list[str], top=1.55, left=0.85, width=6.2, accent_lead=False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}" if not accent_lead else item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK if not accent_lead else MUTED
        p.font.name = "Calibri"
        p.space_after = Pt(14)
        p.level = 0


def add_feature_cards(slide, cards: list[tuple[str, str]], top=1.45) -> None:
    widths = [3.85, 3.85, 3.85]
    for i, (title, body) in enumerate(cards):
        left = Inches(0.75 + i * 4.05)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(top), Inches(widths[i]), Inches(4.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = BOX_BORDER
        tf = shape.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.3)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.font.name = "Georgia"
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED
        p2.font.name = "Calibri"
        p2.space_before = Pt(10)


def add_box(slide, left, top, width, height, text, fill=BOX_FILL, font_size=12, bold=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BOX_BORDER
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = DARK
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = ACCENT
    conn.line.width = Pt(2)
    return conn


def build_workflow_diagram(slide) -> None:
    add_slide_title(slide, "Tech Stack & System Workflow")

    add_box(slide, Inches(5.2), Inches(1.3), Inches(2.6), Inches(0.6), "User Browser", bold=True)
    add_box(
        slide, Inches(3.5), Inches(2.2), Inches(6.1), Inches(0.95),
        "Frontend — Next.js 15 · TypeScript · Tailwind\nTipTap Editor · StPageFlip Reader · TanStack Query",
        bold=True,
    )
    add_box(
        slide, Inches(3.5), Inches(3.55), Inches(6.1), Inches(0.95),
        "Backend — Go (Gin) REST API :8080\nJWT · RBAC · GORM · Zap · Docker Compose",
        bold=True,
    )
    add_box(
        slide, Inches(0.9), Inches(4.95), Inches(3.3), Inches(0.95),
        "PostgreSQL 16\nUsers · Documents · Library · Audit",
        bold=True,
    )
    add_box(
        slide, Inches(4.6), Inches(4.95), Inches(2.8), Inches(0.95),
        "File Storage\nDrafts · DOCX · Covers",
        bold=True,
    )
    add_box(
        slide, Inches(7.9), Inches(4.95), Inches(3.5), Inches(0.95),
        "AI (Optional)\nHugging Face — Qwen 2.5 7B",
        bold=True, fill=AI_FILL,
    )

    add_arrow(slide, Inches(6.5), Inches(1.9), Inches(6.5), Inches(2.2))
    add_arrow(slide, Inches(6.5), Inches(3.15), Inches(6.5), Inches(3.55))
    add_arrow(slide, Inches(5.0), Inches(4.5), Inches(2.55), Inches(4.95))
    add_arrow(slide, Inches(6.5), Inches(4.5), Inches(6.0), Inches(4.95))
    add_arrow(slide, Inches(8.0), Inches(4.5), Inches(9.65), Inches(4.95))

    note = slide.shapes.add_textbox(Inches(0.75), Inches(6.2), Inches(11.8), Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = (
        "Flow: Browser → Next.js Portal → Gin API → PostgreSQL / uploads  |  "
        "AI: proofread, summaries, recommendations via Hugging Face"
    )
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.font.name = "Calibri"


def add_role_table(slide) -> None:
    add_slide_title(slide, "User Role Permissions")

    note = slide.shapes.add_textbox(Inches(0.78), Inches(1.05), Inches(11.5), Inches(0.45))
    p = note.text_frame.paragraphs[0]
    p.text = (
        "EUKOV uses a two-tier model: Readers consume content; one Super Admin account "
        "(bootstrapped at deploy) holds unified Author + Admin + Super Admin powers."
    )
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"
    p.font.italic = True

    rows, cols = 6, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.75), Inches(1.65), Inches(11.85), Inches(4.6))
    table = table_shape.table

    headers = ["Feature Activity", "Reader", "Super Admin (single account)"]
    data = [
        ("Interactive 3D flipbook reader, bookmarks, in-book search, TTS", "✓", "✓"),
        ("Library browse, issue books, reading progress & continue reading", "✓", "✓"),
        ("Docket editor — write, autosave, DOCX import/export, publish", "—", "✓"),
        ("Author review queue, inbox messaging, access-key generation", "—", "✓"),
        ("Audit logs, published-script takedowns, platform governance", "—", "✓"),
    ]

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.name = "Calibri"
            paragraph.font.color.rgb = DARK

    for r, row_data in enumerate(data, start=1):
        for c, value in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                paragraph.font.color.rgb = DARK if c == 0 else ACCENT
                if c > 0:
                    paragraph.alignment = PP_ALIGN.CENTER


def add_roadmap(slide) -> None:
    add_slide_title(slide, "Platform Future Roadmap")

    phases = [
        ("Phase 1", "Host AI on internal servers for faster, offline-ready proofread & summaries."),
        ("Phase 2", "Premium cloud TTS for consistent audiobook-quality narration."),
        ("Phase 3", "Progressive Web App + mobile-first reader redesign."),
        ("Phase 4", "Real-time co-authoring and collaborative writers' rooms."),
    ]
    positions = [1.55, 2.85, 4.15, 5.45]
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.2), Inches(3.55), Inches(10.9), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xF5, 0xC5, 0x6A)
    line.line.fill.background()

    for i, ((phase, desc), x) in enumerate(zip(phases, positions)):
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.42), Inches(0.28), Inches(0.28)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        top = 2.0 if i % 2 == 1 else 4.35
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x - 1.1), Inches(top), Inches(2.5), Inches(1.35),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = BOX_FILL
        box.line.color.rgb = BOX_BORDER
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        p = tf.paragraphs[0]
        p.text = phase
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.name = "Georgia"
        p.font.color.rgb = DARK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.name = "Calibri"
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(6)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1)
    add_center_title_slide(
        s1,
        "EUKOV Management Portal",
        "An intuitive, unified digital publishing platform to write, review, publish, "
        "and explore long-form works in a 3D flipbook reader.",
        "github.com/HarrisSV/Eukov",
    )

    # 2 — Problem
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2)
    add_slide_title(s2, "The Publishing Problem")
    add_two_column_boxes(
        s2,
        "Fragmented Systems",
        "Traditional publishing scatters writing, review, distribution, and reading "
        "across disconnected tools — slowing authors and operators alike.",
        "Lacking Reader UX",
        "Basic portals miss bookmarks, in-book search with highlights, reading-progress "
        "sync, immersive pagination, and accessible audio for long-form archives "
        "(e.g. Project Gutenberg).",
    )

    # 3 — Solution
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3)
    add_slide_title(s3, "Our Unified Solution")
    lead = s3.shapes.add_textbox(Inches(0.78), Inches(1.05), Inches(11), Inches(0.4))
    lead.text_frame.paragraphs[0].text = "Bringing all publishing tasks into one editorial hub:"
    lead.text_frame.paragraphs[0].font.size = Pt(14)
    lead.text_frame.paragraphs[0].font.color.rgb = ACCENT
    lead.text_frame.paragraphs[0].font.name = "Calibri"
    add_bullet_list(
        s3,
        [
            "Universal Docket — create, autosave, and publish manuscripts in a rich TipTap editor",
            "3D Flipbook Reader — StPageFlip animations, fullscreen, chapters, TTS & word highlighting",
            "Two-Tier Roles — Readers explore the library; one Super Admin runs the entire platform "
            "(Author + Admin + Super Admin capabilities unified in a single bootstrap account)",
            "Optional Qwen AI — proofread, quick/full summaries, and personalized recommendations",
        ],
        top=1.55,
        width=11.5,
    )

    # 4 — Features
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4)
    add_slide_title(s4, "Platform Key Features")
    add_feature_cards(
        s4,
        [
            (
                "For Readers",
                "Browse the global library, issue books to your docket, read in fullscreen "
                "with bookmarks, yellow search highlights, progress sync, and browser TTS narration.",
            ),
            (
                "For Super Admin",
                "Single platform operator: write & publish in the docket, review author applications, "
                "generate access keys, manage takedowns, audit logs, and inbox messaging.",
            ),
            (
                "AI Helper (Qwen)",
                "Optional Hugging Face integration for grammar proofread, book digests, "
                "full-book AI summaries, and genre-aware library recommendations.",
            ),
        ],
    )

    # 5 — Tech stack
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5)
    build_workflow_diagram(s5)

    # 6 — Roles
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6)
    add_role_table(s6)

    # 7 — Limitations
    s7 = prs.slides.add_slide(blank)
    set_slide_bg(s7)
    add_slide_title(s7, "Platform Limitations")
    sub = s7.shapes.add_textbox(Inches(0.78), Inches(1.05), Inches(5), Inches(0.35))
    sub.text_frame.paragraphs[0].text = "Current growth constraints:"
    sub.text_frame.paragraphs[0].font.size = Pt(14)
    sub.text_frame.paragraphs[0].font.color.rgb = ACCENT
    sub.text_frame.paragraphs[0].font.name = "Calibri"
    add_bullet_list(
        s7,
        [
            "AI Speed Cap — depends on Hugging Face API latency, rate limits, and token availability",
            "Basic Narration — TTS uses the browser speech engine; voice quality varies by OS",
            "Web Only — desktop-first portal; native mobile apps not yet available",
            "Single Super Admin — platform governance centralized on one bootstrap account",
        ],
        top=1.55,
        width=6.5,
    )

    # 8 — Roadmap
    s8 = prs.slides.add_slide(blank)
    set_slide_bg(s8)
    add_roadmap(s8)

    # 9 — Closing
    s9 = prs.slides.add_slide(blank)
    set_slide_bg(s9)
    glow = s9.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.5), Inches(-0.5), Inches(4), Inches(4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xC4)
    glow.line.fill.background()
    box = s9.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Answers"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "Georgia"
    p.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = "Thank you for exploring the future of digital literature."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(17)
    p2.font.color.rgb = MUTED
    p2.font.name = "Calibri"
    p2.space_before = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = "EUKOV Management Portal  ·  github.com/HarrisSV/Eukov"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(13)
    p3.font.color.rgb = ACCENT
    p3.font.name = "Calibri"
    p3.space_before = Pt(36)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
